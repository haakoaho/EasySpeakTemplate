from dataclasses import dataclass
import datetime
import requests
from typing import List, Dict
from bs4 import BeautifulSoup
import re
from pptx import Presentation

# Your existing dataclasses remain the same.
@dataclass
class TimeSlot:
    time: str
    role: str
    presenter: str
    event: str
    duration_green: str
    duration_amber: str
    duration_red: str


@dataclass
class Speaker:
    position: str
    name: str
    project: str 
    title: str  
    description: str  
    time: str
    duration_green: str
    duration_amber: str
    duration_red: str


@dataclass
class MeetingInfo:
    club_name: str
    district: str
    division: str
    area: str
    club_number: str
    meeting_date: str
    next_meeting_date: str
    meeting_time: str
    venue: str
    schedule: str
    word_of_the_day: str
    meething_theme: str


@dataclass
class ToastmastersMeeting:
    meeting_info: MeetingInfo
    agenda_items: List[TimeSlot]
    speakers: List[Speaker]
    attending_members: List[str]
    next_meeting: str


def _get_day_suffix(day):
    """
    Helper function to get the correct day suffix for a given day of the month.
    """
    if 4 <= day <= 20 or 24 <= day <= 30:
        return "th"
    else:
        return ["st", "nd", "rd"][day % 10 - 1] if day % 10 in [1, 2, 3] else "th"

def next_week(date_string: str) -> str:
    """
    Converts a date string in the format "Weekday dayOfMonth month Year"
    and returns a new date string exactly 7 days in the future.

    Args:
        date_string: The input date string, e.g., "Tuesday 2nd September 2025".

    Returns:
        The new date string, 7 days in the future, in the same format.
    """
    # Use a regex to remove the ordinal suffix (e.g., 'st', 'nd', 'rd', 'th')
    # to allow datetime.strptime to parse the string correctly.
    cleaned_date_string = re.sub(r'(st|nd|rd|th)', '', date_string)

    try:
        # Parse the cleaned string into a datetime object.
        # %A: Full weekday name
        # %d: Day of the month as a zero-padded decimal number
        # %B: Full month name
        # %Y: Year with century as a decimal number
        parsed_date = datetime.datetime.strptime(cleaned_date_string, '%A %d %B %Y')
    except ValueError as e:
        return f"Error: Could not parse the date string. Check the format. Details: {e}"

    # Add exactly 7 days to the parsed date.
    next_week_date = parsed_date + datetime.timedelta(days=7)

    # Get the day of the month from the new date.
    day = next_week_date.day

    # Get the correct suffix for the new day.
    suffix = _get_day_suffix(day)

    # Format the new date string with the correct suffix.
    # We use f-string formatting to insert the day and suffix manually.
    formatted_date = next_week_date.strftime(f'%A {day}{suffix} %B %Y')

    return formatted_date
    # Remove ordinal suffix from the day
    clean_str = (
        date_str.replace("st", "")
        .replace("nd", "")
        .replace("rd", "")
        .replace("th", "")
    )

    # Parse date
    dt = datetime.strptime(clean_str, "%A %d %B %Y")

    # Add 7 days (next Tuesday)
    next_dt = dt + datetime.timedelta(days=7)

    # Format with ordinal
    day_with_suffix = add_ordinal(next_dt.day)
    return next_dt.strftime(f"%A {day_with_suffix} %B %Y")

def scrape_easyspeak_agenda(html_content: str) -> ToastmastersMeeting:
    """
    Scrapes the HTML content of an EasySpeak agenda page and extracts meeting details,
    agenda items, speakers, attending members, and next meeting information.

    Args:
        html_content: The full HTML content of the EasySpeak agenda page.

    Returns:
        A ToastmastersMeeting dataclass object containing the scraped data.
    """
    soup = BeautifulSoup(html_content, 'html.parser')

    # --- Extract Meeting Information ---
    club_name_elem = soup.find('a', class_='maintitle')
    club_name = club_name_elem.text.strip() if club_name_elem else ""

    district_info_elem = soup.find('span', class_='gensmall', string=re.compile(r'District \d+'))
    district = division = area = club_number = ""
    if district_info_elem:
        district_text = district_info_elem.text.strip()
        parts = district_text.split(', ')
        if len(parts) >= 4:
            district = parts[0].replace('District ', '')
            division = parts[1].replace('Division ', '')
            area = parts[2].replace('Area ', '')
            club_number = parts[3].replace('Club Number ', '')

    meeting_date = ""
    next_meeting_date = ""
    word_of_the_day = ""
    postbody_spans = soup.find_all('span', class_='postbody')
    for span in postbody_spans:
        date_match = re.search(r'(Monday|Tuesday|Wednesday|Thursday|Friday|Saturday|Sunday)\s+\d+\w*\s+\w+\s+\d{4}',
                               span.text.strip())
        word_of_the_day_match = re.search(r'(?<=Word of the Day).*?$', span.text)
        if word_of_the_day_match:
            word_of_the_day = word_of_the_day_match.group().strip()
        if date_match:
            meeting_date = date_match.group()
            next_meeting_date = next_week(meeting_date)
            break

    meeting_time_elem = soup.find('b', string=re.compile(r'\d{1,2}:\d{2}'))
    meeting_time = meeting_time_elem.text.strip() if meeting_time_elem else ""

    venue = ""
    for span in postbody_spans:
        if 'Venue ' in span.text:
            venue = span.text.strip().replace('Venue ', '').strip()
            break

    schedule_elem = soup.find('span', class_='gensmall', string=re.compile(r'Every'))
    schedule = schedule_elem.text.strip() if schedule_elem else ""

    meeting_info = MeetingInfo(club_name, district, division, area, club_number, meeting_date,next_meeting_date, meeting_time, venue,
                               schedule, word_of_the_day,"N/A")

    # --- Extract Agenda Items and Speakers ---
    agenda_items = []
    speakers = []
    main_table = soup.find('table', {'border': '0', 'cellpadding': '1', 'cellspacing': '2'})
    if main_table:
        rows = main_table.find_all('tr')
        i = 0
        while i < len(rows):
            row = rows[i]
            cells = row.find_all('td')

            # Check if this row is a main agenda item row (e.g., has a time and enough cells)
            time_span_in_first_cell = cells[0].find('span', class_='gensmall') if cells and len(cells) > 0 else None

            # A valid agenda row should have at least 5 cells.
            # Time is optional — it may be blank if grouped under a previous item (like for multiple evaluations).
            if len(cells) < 5:
                i += 1
                continue

            # Get the time if available
            time_slot = ""
            if time_span_in_first_cell and re.match(r'\d{1,2}:\d{2}', time_span_in_first_cell.text.strip()):
                time_slot = time_span_in_first_cell.text.strip()
            else:
                # Inherit time from previous agenda item if blank
                time_slot = agenda_items[-1].time if agenda_items else "TBA"
            role = cells[1].find('span', class_='gen').text.strip() if cells[1].find('span', class_='gen') else ""
            presenter = cells[2].find('span', class_='gen').text.strip() if cells[2].find('span', class_='gen') else ""

            # The speech title is in the 'event' column for Speakers
            event = cells[3].find('span', class_='gensmall').text.strip() if cells[3].find('span',
                                                                                           class_='gensmall') else ""

            duration_green, duration_amber, duration_red = "", "", ""
            duration_span_elem = cells[4].find('span', class_='gensmall')
            if duration_span_elem:
                duration_parts = re.split(r'\s+', duration_span_elem.text.strip())
                if len(duration_parts) > 0: duration_green = duration_parts[0]
                if len(duration_parts) > 1: duration_amber = duration_parts[1]
                if len(duration_parts) > 2: duration_red = duration_parts[2]

            # Add to general agenda items
            agenda_items.append(
                TimeSlot(time_slot, role, presenter, event, duration_green, duration_amber, duration_red))

            # If it's a speaker, extract project and detailed description from the next row
            if "Speaker" in role:
                project = "TBA"  
                title = event  
                description = ""

                speaker_detail_row = None

                # Iterate through the next siblings to find the specific detail row
                # The detail row should contain a <td> with colspan="3" AND align="left"
                potential_next_rows = rows[i + 1:]  # Get all rows after the current one

                for potential_row in potential_next_rows:
                    target_td = potential_row.find('td', {'colspan': '3', 'align': 'left'})

                    if target_td:
                        speaker_detail_row = potential_row
                        break  

                if speaker_detail_row:
                    # The project/description is in the <td> with colspan="3" and align="left"
                    project_desc_td = speaker_detail_row.find('td', {'colspan': '3', 'align': 'left'})


                    if project_desc_td:
                        # Find the span directly within this td
                        project_desc_span = project_desc_td.find('span', class_='gensmall', valign="top")

                        if project_desc_span:
                            # The project line itself is inside an <i> tag (e.g., Pathways project)
                            i_tag = project_desc_span.find('i')


                            if i_tag:
                                full_project_line = i_tag.text.strip()

                                # The 'project' is the part before ' - ' in the i_tag (if applicable)
                                # The 'description' is the rest of the text in the span, after the i_tag.
                                project_parts = full_project_line.split(' - ', 1)
                                project = project_parts[0].strip()

                                # The description includes all subsequent text nodes within the span, separated by <br>
                                desc_lines = []
                                # Add the part after ' - ' from the i_tag to the description if it's there
                                if len(project_parts) > 1:
                                    # This is usually the specific title *from the Pathways manual* for the speech type
                                    # or a short description of the Pathways objective.
                                    desc_lines.append(project_parts[1].strip())

                                # Combine all description parts, filtering out empty strings
                                description = " ".join([line for line in desc_lines if line]).strip()

                            else:  # If no <i> tag is found within the span (e.g., a custom speech not tied to Pathways)
                                all_strings_in_span = [s.strip() for s in project_desc_span.stripped_strings if
                                                       s.strip()]
                                if all_strings_in_span:
                                    project = "N/A (No Pathways Info)"  # Indicate no Pathways info
                                    description = " ".join(
                                        all_strings_in_span).strip()  # All text in span is description
                                else:
                                    project = "N/A (No Pathways Info)"
                                    description = ""
                        else:
                            project = "TBA"
                            description = ""
                    else:
                        project = "TBA"
                        description = ""
                else:  # No valid speaker_detail_row found at all
                    project = "TBA"
                    description = ""

                # Add the speaker information to the speakers list
                speakers.append(
                    Speaker(role, presenter, project, title, description, time_slot, duration_green, duration_amber,
                            duration_red))

                # Advance 'i' past the speaker detail row if it was found
                if speaker_detail_row:
                    try:
                        # Find the index of the found speaker_detail_row in the original 'rows' list
                        # This ensures 'i' correctly jumps past the detail row, no matter how many
                        # "filler" rows were in between.
                        detail_row_index = rows.index(speaker_detail_row, i + 1)
                        i = detail_row_index  # Set 'i' to the index of the detail row.
                    except ValueError:
                        # This should ideally not happen if speaker_detail_row was successfully found and is in `rows`.
                        pass

            i += 1  # Always increment for the current main agenda row, or the detail row if it was just processed.

    # --- Extract Attending Members ---
    attending_members = []
    attending_section = soup.find('span', class_='cattitle', string='Attending')
    if attending_section and attending_section.find_parent('table'):
        # Find the <tr> that contains the 'Attending' heading
        attending_header_row = attending_section.find_parent('tr')
        if attending_header_row:
            # The actual member list is usually in the next <tr> sibling
            members_row = attending_header_row.find_next_sibling('tr')
            if members_row:
                member_cells = members_row.find_all('td', class_='gensmall')
                for cell in member_cells:
                    raw_names = re.split(r'[,;\n\r]+', cell.text.strip())
                    for name in raw_names:
                        cleaned_name = name.strip()
                        if cleaned_name and cleaned_name != "Member":  # Exclude "Member" if it appears as a placeholder
                            attending_members.append(cleaned_name)

    # --- Extract Next Meeting Info ---
    next_meeting = ""
    next_meeting_elem = soup.find('span', class_='cattitle', string='Next Meeting')
    if next_meeting_elem and next_meeting_elem.find_next_sibling('span', class_='gensmall'):
        next_meeting = next_meeting_elem.find_next_sibling('span', class_='gensmall').text.strip()

    return ToastmastersMeeting(meeting_info, agenda_items, speakers, attending_members, next_meeting)


def get_roles_and_presenters(meeting_data: ToastmastersMeeting) -> Dict[str, Dict[str, str]]:
    """
    Extracts a dictionary mapping roles to their presenters and time frames.
    """
    roles_info = {}
    for item in meeting_data.agenda_items:
        if not item.role or item.role == 'Break':
            continue
        # Sanitize role name to be a valid placeholder key
        sanitized_role = re.sub(r'[^a-zA-Z0-9]', '', item.role)
        roles_info[sanitized_role] = {
            'presenter': item.presenter if item.presenter else "N/A",
            'min_time': item.duration_green if item.duration_green else "N/A",
            'max_time': item.duration_red if item.duration_red else "N/A"
        }
    return roles_info

def update_pptx_presentation(template_path: str, output_path: str, meeting_data: ToastmastersMeeting):
    """
    Updates a PowerPoint (.pptx) presentation by replacing placeholders.
    Args:
        template_path: Path to the template .pptx file.
        output_path: Path to save the updated .pptx file.
        meeting_data: The ToastmastersMeeting dataclass object containing all scraped data.
    """
    prs = Presentation(template_path)

    # Prepare replacements
    replacements = {}

    replacements["{{meeting_date}}"] = meeting_data.meeting_info.meeting_date or "N/A"
    replacements["{{next_meeting_date}}"] = meeting_data.meeting_info.next_meeting_date or "N/A"
    replacements["{{word_of_the_day}}"] = meeting_data.meeting_info.word_of_the_day or "N/A"
    

    roles_presenters_time = get_roles_and_presenters(meeting_data)
    for role_key, info in roles_presenters_time.items():
        replacements[f"{{{{{role_key}_Presenter}}}}"] = info.get("presenter", "N/A")
        replacements[f"{{{{{role_key}_Min}}}}"] = info.get("min_time", "N/A")
        replacements[f"{{{{{role_key}_Max}}}}"] = info.get("max_time", "N/A")

    for i, speaker in enumerate(meeting_data.speakers):
        if i < 3:
            num = i + 1
            replacements[f"{{{{Speaker{num}_Name}}}}"] = speaker.name or "N/A"
            replacements[f"{{{{Speaker{num}_Project}}}}"] = speaker.project or "N/A"
            replacements[f"{{{{Speaker{num}_Title}}}}"] = speaker.title or "N/A"
            replacements[f"{{{{Speaker{num}_Description}}}}"] = speaker.description or "N/A"
            replacements[f"{{{{Speaker{num}_Min}}}}"] = speaker.duration_green or "N/A"
            replacements[f"{{{{Speaker{num}_Max}}}}"] = speaker.duration_red or "N/A"

    # Replace text in all shapes across all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)

    prs.save(output_path)
    print(f"\nPPTX presentation successfully created at: {output_path}")


def update_forms(meeting_data: ToastmastersMeeting):
    feedback_form = "https://script.google.com/macros/s/AKfycbz7CpmaWJ3K0A3JQydu9F2z-h1lPNwkh7OeGl-ia5stQ-XJY6i7CwtOS6Sv9959e-jJ/exec"
    speaker_form = "https://script.google.com/macros/s/AKfycbxcFD8uWhxh0wg6ZXkGU7ZmqUpCTtzTklkcV6h1JamJ0gSX--z7jvk6WtfPdTGElM1E/exec"
    evaluator_form = "https://script.google.com/macros/s/AKfycbydCQEm60REI0gy7SC1g4rpQU3hlOwI-FUwGwE2GqUiQamGVY0N3QTua5GQOh-lOdLU/exec"
    speakers = []
    for speak in meeting_data.speakers:
        speakers.append(speak.name)
    evaluators = []
    for item in meeting_data.agenda_items:
        if "Evaluate speech" or "Table Topics Evaluator" in item.event:
            evaluators.append(item.presenter)
    res = requests.post(feedback_form, json={"options": speakers})
    res = requests.post(speaker_form, json={"options": speakers})
    res = requests.post(evaluator_form, json={"options": evaluators})


if __name__ == "__main__":

    meething_theme = input("Enter the meeting theme (or leave blank): ").strip()
    print("=== EasySpeak Agenda Parser ===\n")

    print("➡️  1. Open the agenda page in your browser.")
    print("➡️  2. Open Developer Tools (F12), go to the Console tab.")
    print("➡️  3. Run this command to copy the full HTML:\n")
    print("       copy(document.documentElement.outerHTML);\n")
    print("➡️  4. Paste the copied HTML below (right-click or Ctrl+V).")
    print("       When you're done, type `END` on a new line and press Enter.\n")

    # Read multiline HTML input until 'END' is typed
    html_lines = []
    while True:
        line = input()
        if line.strip() == "END":
            break
        html_lines.append(line)

    html_content = "\n".join(html_lines)

    if not html_content.strip():
        print("❌ No HTML content provided. Exiting.")
        exit(1)

    # Hardcoded file names
    template_file = "template.pptx"
    output_file = "presentation.pptx"

    try:
        agenda_object = scrape_easyspeak_agenda(html_content)
        agenda_object.meeting_info.meething_theme = meething_theme if meething_theme else "N/A"

        # Print role assignments
        roles_presenters_time = get_roles_and_presenters(agenda_object)
        print("\n--- Roles, Presenters, and Time Frames ---")
        for role, info in roles_presenters_time.items():
            print(f"{role}: {info['presenter']} (Min: {info['min_time']}, Max: {info['max_time']})")

        # Print speaker data
        print("\n--- Speaker Information ---")
        for i, speaker in enumerate(agenda_object.speakers):
            if i < 3:
                print(
                    f"Speaker {i + 1}: {speaker.name}, Project: '{speaker.project}', Title: '{speaker.title}', Description: '{speaker.description}'"
                )

        print(f"\nMeeting Date: {agenda_object.meeting_info.meeting_date}")
        print(f"Next Meeting Date: {agenda_object.meeting_info.next_meeting_date}")
        print(f"Word of the Day: {agenda_object.meeting_info.word_of_the_day}")

        print(f"\n📤 Generating presentation from '{template_file}'...")
        update_pptx_presentation(template_file, output_file, agenda_object)

        print("\n🔗 Updating Google Forms...")
        update_forms(agenda_object)

    except Exception as e:
        print(f"\n❌ An error occurred: {e}")
